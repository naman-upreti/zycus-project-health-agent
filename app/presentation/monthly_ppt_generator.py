from datetime import datetime
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.presentation.ppt_generator import (
    NAVY,
    NAVY_SOFT,
    NAVY_CARD,
    INK,
    MUTED,
    ICE,
    WHITE,
    LIGHT_GRAY,
    CARD_BORDER,
    GREEN,
    AMBER,
    RED,
    GREY_STATUS,
    HEADING_FONT,
    BODY_FONT,
    SLIDE_W,
    SLIDE_H,
    MARGIN,
)


class MonthlyPPTGenerator:

    def __init__(self, reports):
        """
        reports: list of dicts, each shaped as

            {
                "metrics": ProjectMetrics,
                "health": HealthResult,
                "summary": str,
            }
        """
        self.reports = reports or []

    # -----------------------------------------------------------------
    # Generic, defensive accessors (same pattern as PPTGenerator)
    # -----------------------------------------------------------------

    @staticmethod
    def _get(obj, name, default=None):
        if obj is None:
            return default
        if isinstance(obj, dict):
            return obj.get(name, default)
        return getattr(obj, name, default)

    def _rag_color(self, status):
        status = (status or "").strip().lower()
        return {"green": GREEN, "amber": AMBER, "red": RED}.get(status, GREY_STATUS)

    def _score_color(self, value):
        try:
            value = float(value)
        except (TypeError, ValueError):
            return GREY_STATUS
        if value >= 75:
            return GREEN
        if value >= 50:
            return AMBER
        return RED

    def _fmt_score(self, value, suffix="%"):
        try:
            return f"{float(value):.0f}{suffix}"
        except (TypeError, ValueError):
            return "N/A"

    def _normalize_recommendation(self, item):
        if isinstance(item, dict):
            text = item.get("text") or item.get("action") or item.get("recommendation") or str(item)
            return text
        if isinstance(item, str):
            return item
        text = self._get(item, "text") or self._get(item, "action") or str(item)
        return text

    # -----------------------------------------------------------------
    # Portfolio-level aggregation helpers
    # -----------------------------------------------------------------

    def _report_project_name(self, report, fallback):
        return self._get(report.get("metrics"), "project_name", fallback)

    def _report_overall_score(self, report):
        try:
            return float(self._get(report.get("health"), "overall_score", 0) or 0)
        except (TypeError, ValueError):
            return 0.0

    def _report_rag_status(self, report):
        return str(self._get(report.get("health"), "rag_status", "N/A"))

    def _rag_counts(self):
        counts = Counter()
        for report in self.reports:
            status = self._report_rag_status(report).strip().lower()
            counts[status] += 1
        return (
            counts.get("green", 0),
            counts.get("amber", 0),
            counts.get("red", 0),
        )

    def _average_portfolio_score(self):
        if not self.reports:
            return 0.0
        scores = [self._report_overall_score(r) for r in self.reports]
        return sum(scores) / len(scores)

    def _portfolio_health_label(self, average_score):
        if average_score >= 75:
            return "Green", GREEN
        if average_score >= 50:
            return "Amber", AMBER
        return "Red", RED

    def _combined_executive_summary(self):
        """Combine every per-project 'summary' string into ONE concise
        executive-level narrative, without simply printing each summary
        separately. Since this is presentation-layer only (no LLM call,
        no business logic), the combination is done by distilling the
        lead sentence of each project's summary into a single portfolio
        narrative paragraph."""
        summaries = [
            (self._report_project_name(r, "Unnamed Project"), (r.get("summary") or "").strip())
            for r in self.reports
        ]
        summaries = [(name, s) for name, s in summaries if s]

        if not summaries:
            return "No executive summaries were available for this reporting cycle."

        lead_sentences = []
        for name, summary in summaries:
            first_sentence = summary.split(".")[0].strip()
            if first_sentence:
                lead_sentences.append(f"{name} — {first_sentence}.")

        green, amber, red = self._rag_counts()
        avg_score = self._average_portfolio_score()
        health_label, _ = self._portfolio_health_label(avg_score)

        intro = (
            f"Across {len(self.reports)} project(s) reviewed this cycle, the portfolio "
            f"stands at an average health score of {avg_score:.0f}/100 ({health_label} "
            f"overall), with {green} project(s) Green, {amber} Amber, and {red} Red."
        )

        narrative = intro + " " + " ".join(lead_sentences)
        return narrative

    def _combined_recommendations(self):
        """Collect every recommendation across all reports, de-duplicate
        while preserving first-seen order, and return the list."""
        seen = set()
        combined = []
        for report in self.reports:
            recs = self._get(report.get("health"), "recommendations", []) or []
            for rec in recs:
                text = self._normalize_recommendation(rec).strip()
                if not text:
                    continue
                key = text.lower()
                if key in seen:
                    continue
                seen.add(key)
                combined.append(text)
        if not combined:
            combined = ["No specific recommendations were raised for this cycle."]
        return combined

    def _emerging_risk_insights(self):
        """Portfolio-level trend insights rather than a per-project list."""
        total_critical = 0
        total_on_hold = 0
        behind_schedule = 0
        milestone_delays = 0

        for report in self.reports:
            metrics = report.get("metrics")
            health = report.get("health")

            total_critical += int(self._get(metrics, "critical_tasks", 0) or 0)
            total_on_hold += int(self._get(metrics, "on_hold_tasks", 0) or 0)

            schedule_score = self._get(health, "schedule_score", None)
            if schedule_score is not None:
                try:
                    if float(schedule_score) < 50:
                        behind_schedule += 1
                except (TypeError, ValueError):
                    pass

            milestone_score = self._get(health, "milestone_score", None)
            if milestone_score is not None:
                try:
                    if float(milestone_score) < 50:
                        milestone_delays += 1
                except (TypeError, ValueError):
                    pass

        return [
            ("Total Critical Tasks", str(total_critical), NAVY),
            ("Total On-Hold Tasks", str(total_on_hold), AMBER),
            ("Projects Behind Schedule", str(behind_schedule), RED),
            ("Projects With Milestone Delays", str(milestone_delays), RED),
        ]

    # -----------------------------------------------------------------
    # Low-level drawing helpers (same visual language as PPTGenerator)
    # -----------------------------------------------------------------

    def _blank_slide(self, prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _fill_background(self, slide, color):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def _textbox(self, slide, x, y, w, h, wrap=True):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = wrap
        return box, tf

    def _style_run(self, paragraph, text=None, font=BODY_FONT, size=14, bold=False,
                    italic=False, color=INK, align=None, space_after=None, space_before=None):
        if text is not None:
            paragraph.text = text
        paragraph.font.name = font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        paragraph.font.color.rgb = color
        if align is not None:
            paragraph.alignment = align
        if space_after is not None:
            paragraph.space_after = Pt(space_after)
        if space_before is not None:
            paragraph.space_before = Pt(space_before)
        return paragraph

    def _rounded_card(self, slide, x, y, w, h, fill=LIGHT_GRAY, line_color=CARD_BORDER, line_pt=0.75, radius=None):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        if line_color is None:
            card.line.fill.background()
        else:
            card.line.color.rgb = line_color
            card.line.width = Pt(line_pt)
        card.shadow.inherit = False
        if radius is not None:
            try:
                card.adjustments[0] = radius
            except (IndexError, ValueError):
                pass
        return card

    def _dot(self, slide, x, y, d, color):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.shadow.inherit = False
        return dot

    def _pill(self, slide, x, y, w, h, text, fill, text_color=WHITE, size=16):
        pill = self._rounded_card(slide, x, y, w, h, fill=fill, line_color=None, radius=0.5)
        tf = pill.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        p = tf.paragraphs[0]
        self._style_run(p, text, font=BODY_FONT, size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        tf.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE
        return pill

    def _section_title(self, slide, x, y, w, text, color=NAVY, size=18):
        _, tf = self._textbox(slide, x, y, w, Inches(0.4))
        self._style_run(tf.paragraphs[0], text, font=HEADING_FONT, size=size, bold=True, color=color)
        return tf

    def _bullet_list(self, slide, x, y, w, h, items, color=INK, size=13.5, space_after=6):
        _, tf = self._textbox(slide, x, y, w, h)
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            self._style_run(p, f"•  {item}", font=BODY_FONT, size=size, color=color, space_after=space_after)
        return tf

    def _title_bar(self, slide, kicker, title, page_no, section_color=NAVY):
        _, ktf = self._textbox(slide, MARGIN, Inches(0.4), Inches(9), Inches(0.3))
        self._style_run(ktf.paragraphs[0], kicker.upper(), font=BODY_FONT, size=11,
                        bold=True, color=MUTED)

        _, ttf = self._textbox(slide, MARGIN, Inches(0.68), Inches(11), Inches(0.65))
        self._style_run(ttf.paragraphs[0], title, font=HEADING_FONT, size=27, bold=True, color=section_color)

        _, ptf = self._textbox(slide, Inches(12.2), Inches(0.42), Inches(0.6), Inches(0.3))
        self._style_run(ptf.paragraphs[0], f"{page_no:02d}", font=BODY_FONT, size=11,
                        color=MUTED, align=PP_ALIGN.RIGHT)

    def _footer(self, slide, page_no, dark=False):
        color = ICE if dark else MUTED
        _, tf = self._textbox(slide, MARGIN, Inches(7.12), Inches(9), Inches(0.3))
        self._style_run(
            tf.paragraphs[0],
            "CONFIDENTIAL  |  Project Health Reporting Agent — Monthly Portfolio Report",
            font=BODY_FONT, size=9, color=color,
        )
        _, ptf = self._textbox(slide, Inches(12.2), Inches(7.12), Inches(0.6), Inches(0.3))
        self._style_run(ptf.paragraphs[0], f"{page_no:02d}", font=BODY_FONT, size=9,
                        color=color, align=PP_ALIGN.RIGHT)

    def _kpi_card(self, slide, x, y, w, h, label, value, score_for_color=None):
        card = self._rounded_card(slide, x, y, w, h, fill=WHITE, line_color=CARD_BORDER, radius=0.10)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.14)

        dot_color = self._score_color(score_for_color) if score_for_color is not None else NAVY
        self._dot(slide, x + Inches(0.15), y + Inches(0.16), Inches(0.14), dot_color)

        label_p = tf.paragraphs[0]
        self._style_run(label_p, "   " + label, font=BODY_FONT, size=12.5, bold=True, color=MUTED)

        value_p = tf.add_paragraph()
        self._style_run(value_p, str(value), font=HEADING_FONT, size=30, bold=True, color=NAVY, space_before=6)

    def _insight_card(self, slide, x, y, w, h, label, value, accent):
        card = self._rounded_card(slide, x, y, w, h, fill=WHITE, line_color=CARD_BORDER, radius=0.10)
        self._dot(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.16), accent)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.margin_top = Inches(0.12)
        label_p = tf.paragraphs[0]
        self._style_run(label_p, label, font=BODY_FONT, size=12.5, bold=True, color=MUTED)
        value_p = tf.add_paragraph()
        self._style_run(value_p, value, font=HEADING_FONT, size=32, bold=True, color=NAVY, space_before=6)

    def _table(self, slide, x, y, w, h, headers, rows, col_widths=None,
            header_fill=NAVY, header_font_color=WHITE, row_fill_even=WHITE,
            row_fill_odd=LIGHT_GRAY, cell_font_size=13, row_colors=None):
        """Readable, consistently-styled table matching the deck's design
        system (navy header, alternating light rows, Calibri body)."""
        n_rows = len(rows) + 1
        n_cols = len(headers)

        graphic_frame = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
        table = graphic_frame.table

        if col_widths:
            for idx, cw in enumerate(col_widths):
                table.columns[idx].width = cw

        # Header row
        for c, header_text in enumerate(headers):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.vertical_anchor = 1  # MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            self._style_run(tf.paragraphs[0], header_text, font=BODY_FONT, size=13, bold=True, color=header_font_color, align=PP_ALIGN.CENTER)

        # Body rows
        for r, row_values in enumerate(rows, start=1):
            fill = row_fill_even if r % 2 == 1 else row_fill_odd
            for c, cell_value in enumerate(row_values):
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
                cell.margin_left = Inches(0.1)
                cell.margin_right = Inches(0.1)
                cell.vertical_anchor = 1  # MIDDLE
                tf = cell.text_frame
                tf.word_wrap = True
                text_color = INK
                if row_colors and row_colors[r - 1] and c == len(row_values) - 1:
                    text_color = row_colors[r - 1]
                self._style_run(
                    tf.paragraphs[0], str(cell_value), font=BODY_FONT, size=cell_font_size,
                    bold=(row_colors is not None and c == len(row_values) - 1),
                    color=text_color, align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT,
                )

        return graphic_frame

    # -----------------------------------------------------------------
    # Slides
    # -----------------------------------------------------------------

    def _slide_portfolio_overview(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, NAVY)

        glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.6), Inches(-2.0), Inches(6.0), Inches(6.0))
        glow.fill.solid()
        glow.fill.fore_color.rgb = NAVY_SOFT
        glow.line.fill.background()
        glow.shadow.inherit = False

        _, ktf = self._textbox(slide, MARGIN, Inches(0.7), Inches(9), Inches(0.35))
        self._style_run(ktf.paragraphs[0], "MONTHLY PORTFOLIO REVIEW", font=BODY_FONT, size=13, bold=True, color=ICE)

        _, ttf = self._textbox(slide, MARGIN, Inches(1.05), Inches(10.5), Inches(1.0))
        self._style_run(ttf.paragraphs[0], "Monthly Portfolio Overview",font=HEADING_FONT, size=34, bold=True, color=WHITE)

        generated_date = datetime.now().strftime("%d %B %Y")
        _, mtf = self._textbox(slide, MARGIN, Inches(1.95), Inches(9), Inches(0.4))
        self._style_run(
            mtf.paragraphs[0], f"Generated: {generated_date}",
            font=BODY_FONT, size=13, color=RGBColor(0x9A, 0xA6, 0xD6),
        )

        projects_reviewed = len(self.reports)
        avg_score = self._average_portfolio_score()
        green, amber, red = self._rag_counts()

        stats = [
            ("PROJECTS REVIEWED", str(projects_reviewed)),
            ("AVERAGE PORTFOLIO SCORE", f"{avg_score:.0f}/100"),
            ("GREEN PROJECTS", str(green)),
            ("AMBER PROJECTS", str(amber)),
            ("RED PROJECTS", str(red)),
        ]

        card_w = Inches(2.35)
        card_h = Inches(2.0)
        gap = Inches(0.2)
        start_x = MARGIN
        start_y = Inches(2.9)
        accent_colors = [ICE, ICE, GREEN, AMBER, RED]
        for i, (label, value) in enumerate(stats):
            x = start_x + Emu(int(i * (card_w + gap)))
            card = self._rounded_card(slide, x, start_y, card_w, card_h, fill=NAVY_CARD,line_color=None, radius=0.08)
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.22)
            p = tf.paragraphs[0]
            self._style_run(p, label, font=BODY_FONT, size=11.5, bold=True, color=ICE)
            vp = tf.add_paragraph()
            self._style_run(vp, value, font=HEADING_FONT, size=30, bold=True,color=accent_colors[i] if i >= 2 else WHITE, space_before=10)

        self._footer(slide, page_no=1, dark=True)
        return slide

    def _slide_project_comparison(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Portfolio Review", "Project Comparison", page_no=2)

        headers = ["Project Name", "Overall Score", "RAG Status"]
        rows = []
        row_colors = []
        for i, report in enumerate(self.reports):
            name = self._report_project_name(report, f"Project {i + 1}")
            score = self._report_overall_score(report)
            status = self._report_rag_status(report)
            rows.append([name, f"{score:.0f}/100", status.upper()])
            row_colors.append(self._rag_color(status))

        if not rows:
            rows = [["No projects available", "N/A", "N/A"]]
            row_colors = [None]

        col_widths = [Inches(7.0), Inches(2.6), Inches(2.6)]
        self._table(
            slide, MARGIN, Inches(1.6), Inches(12.2), Inches(5.0),
            headers=headers, rows=rows, col_widths=col_widths, row_colors=row_colors,
        )

        self._footer(slide, page_no=2)
        return slide

    def _slide_kpi_comparison(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Portfolio Review", "Portfolio KPI Comparison", page_no=3)

        headers = ["Project", "Completion Score", "Schedule Score", "Milestone Score", "Blocker Score"]
        rows = []
        for i, report in enumerate(self.reports):
            name = self._report_project_name(report, f"Project {i + 1}")
            health = report.get("health")
            rows.append([
                name,
                self._fmt_score(self._get(health, "completion_score", None)),
                self._fmt_score(self._get(health, "schedule_score", None)),
                self._fmt_score(self._get(health, "milestone_score", None)),
                self._fmt_score(self._get(health, "blocker_score", None), suffix=""),
            ])

        if not rows:
            rows = [["No projects available", "N/A", "N/A", "N/A", "N/A"]]

        col_widths = [Inches(4.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
        self._table(
            slide, MARGIN, Inches(1.6), Inches(12.2), Inches(5.0),
            headers=headers, rows=rows, col_widths=col_widths, cell_font_size=12.5,
        )

        self._footer(slide, page_no=3)
        return slide

    def _slide_emerging_risks(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Portfolio Review", "Emerging Risks", page_no=4)

        insights = self._emerging_risk_insights()

        card_w = Inches(5.9)
        card_h = Inches(2.0)
        gap_x = Inches(0.4)
        gap_y = Inches(0.3)
        start_x = MARGIN
        start_y = Inches(1.7)
        for i, (label, value, accent) in enumerate(insights):
            col = i % 2
            row = i // 2
            x = start_x + Emu(int(col * (card_w + gap_x)))
            y = start_y + Emu(int(row * (card_h + gap_y)))
            self._insight_card(slide, x, y, card_w, card_h, label, value, accent)

        # Trend narrative card at the bottom
        green, amber, red = self._rag_counts()
        narrative_items = []
        if red > 0:
            narrative_items.append(f"{red} project(s) currently sit in the Red zone and require immediate escalation.")
        if amber > 0:
            narrative_items.append(f"{amber} project(s) are Amber and trending toward risk if unaddressed.")
        if not narrative_items:
            narrative_items.append("No projects currently require escalation; portfolio risk remains contained.")

        y_bottom = Inches(5.85)
        self._rounded_card(slide, MARGIN, y_bottom, Inches(12.2), Inches(1.1), fill=LIGHT_GRAY,
                            line_color=CARD_BORDER, radius=0.08)
        self._bullet_list(slide, MARGIN + Inches(0.3), y_bottom + Inches(0.15), Inches(11.6), Inches(0.85),
                           items=narrative_items, size=13, space_after=4)

        self._footer(slide, page_no=4)
        return slide

    def _slide_executive_summary(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Portfolio Review", "Executive Portfolio Summary", page_no=5)

        narrative = self._combined_executive_summary()

        self._rounded_card(slide, MARGIN, Inches(1.6), Inches(12.2), Inches(5.0), fill=LIGHT_GRAY,
                            line_color=CARD_BORDER, radius=0.05)
        self._section_title(slide, MARGIN + Inches(0.35), Inches(1.85), Inches(11.5), "Summary", size=16)
        _, tf = self._textbox(slide, MARGIN + Inches(0.35), Inches(2.35), Inches(11.5), Inches(4.0))
        self._style_run(tf.paragraphs[0], narrative, font=BODY_FONT, size=14, color=INK)

        self._footer(slide, page_no=5)
        return slide

    def _slide_recommendations(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Portfolio Review", "Executive Recommendations", page_no=6)

        recommendations = self._combined_recommendations()

        top_y = Inches(1.6)
        row_h = Inches(0.85)
        gap = Inches(0.15)
        max_rows = 5
        for i, text in enumerate(recommendations[:max_rows]):
            y = top_y + Emu(int(i * (row_h + gap)))
            self._rounded_card(slide, MARGIN, y, Inches(12.23), row_h, fill=LIGHT_GRAY,
                                line_color=CARD_BORDER, radius=0.14)

            idx_d = Inches(0.5)
            idx_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, MARGIN + Inches(0.2),
                                                y + Inches(0.18), idx_d, idx_d)
            idx_shape.fill.solid()
            idx_shape.fill.fore_color.rgb = NAVY
            idx_shape.line.fill.background()
            idx_shape.shadow.inherit = False
            itf = idx_shape.text_frame
            itf.word_wrap = True
            self._style_run(itf.paragraphs[0], str(i + 1), font=HEADING_FONT, size=16,bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            _, ttf = self._textbox(slide, MARGIN + Inches(0.95), y + Inches(0.18), Inches(11.0), Inches(0.5))
            self._style_run(ttf.paragraphs[0], text, font=BODY_FONT, size=13.5, bold=True, color=INK)

        self._footer(slide, page_no=6)
        return slide

    def _slide_conclusion(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, NAVY)

        _, ktf = self._textbox(slide, MARGIN, Inches(0.6), Inches(9), Inches(0.3))
        self._style_run(ktf.paragraphs[0], "PORTFOLIO CONCLUSION", font=BODY_FONT, size=12,bold=True, color=ICE)

        _, ttf = self._textbox(slide, MARGIN, Inches(0.9), Inches(10), Inches(0.65))
        self._style_run(ttf.paragraphs[0], "Conclusion", font=HEADING_FONT, size=30, bold=True, color=WHITE)

        projects_reviewed = len(self.reports)
        avg_score = self._average_portfolio_score()
        green, amber, red = self._rag_counts()
        health_label, health_color = self._portfolio_health_label(avg_score)

        stats = [
            ("Projects Reviewed", str(projects_reviewed)),
            ("Average Portfolio Score", f"{avg_score:.0f}/100"),
            ("Green", str(green)),
            ("Amber", str(amber)),
            ("Red", str(red)),
        ]
        chip_w = Inches(2.3)
        chip_h = Inches(1.3)
        gap = Inches(0.2)
        start_x = MARGIN
        start_y = Inches(1.9)
        for i, (label, value) in enumerate(stats):
            x = start_x + Emu(int(i * (chip_w + gap)))
            chip = self._rounded_card(slide, x, start_y, chip_w, chip_h, fill=NAVY_CARD,line_color=None, radius=0.1)
            tf = chip.text_frame
            tf.margin_left = Inches(0.15)
            tf.margin_top = Inches(0.14)
            self._style_run(tf.paragraphs[0], label, font=BODY_FONT, size=11.5, bold=True, color=ICE)
            vp = tf.add_paragraph()
            self._style_run(vp, value, font=HEADING_FONT, size=24, bold=True, color=WHITE, space_before=6)

        # Overall Portfolio Health — large final badge
        self._section_title(slide, MARGIN, Inches(3.7), Inches(6), "Overall Portfolio Health",color=WHITE, size=18)
        self._pill(slide, MARGIN, Inches(4.2), Inches(4.0), Inches(1.1), health_label.upper(), health_color, size=28)

        _, note_tf = self._textbox(slide, MARGIN, Inches(5.6), Inches(10.5), Inches(1.0))
        self._style_run(
            note_tf.paragraphs[0],
            "This monthly portfolio review consolidates individual project health "
            "assessments into a single executive view for leadership decision-making.",
            font=BODY_FONT, size=12.5, italic=True, color=RGBColor(0x9A, 0xA6, 0xD6),
        )

        _, thank_tf = self._textbox(slide, MARGIN, Inches(6.6), Inches(7), Inches(0.6))
        self._style_run(thank_tf.paragraphs[0], "Thank You", font=HEADING_FONT, size=24, bold=True, color=WHITE)

        self._footer(slide, page_no=7, dark=True)
        return slide

    # -----------------------------------------------------------------
    # Entry point
    # -----------------------------------------------------------------

    def generate(self, output_path):
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        self._slide_portfolio_overview(prs)
        self._slide_project_comparison(prs)
        self._slide_kpi_comparison(prs)
        self._slide_emerging_risks(prs)
        self._slide_executive_summary(prs)
        self._slide_recommendations(prs)
        self._slide_conclusion(prs)

        prs.save(output_path)