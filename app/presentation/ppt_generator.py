"""
app/presentation/ppt_generator.py

Presentation-layer only: renders a consulting-grade executive PowerPoint from
the metrics / health / summary objects produced upstream by the pipeline
(ExcelParser -> MetricsExtractor -> HealthScoringEngine -> Groq LLM summary).

No business logic lives here. This module only reads attributes off the
objects it receives (defensively, via _get()) and lays them out visually.
"""

from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

from app.models.project_metrics import ProjectMetrics
from app.models.health_result import HealthResult


# =====================================================================
# Design system — one palette / type scale, reused by every slide so the
# deck reads as a single designed artifact rather than six separate ones.
# =====================================================================

NAVY = RGBColor(0x0F, 0x1F, 0x3D)
NAVY_SOFT = RGBColor(0x1D, 0x35, 0x60)
NAVY_CARD = RGBColor(0x25, 0x3F, 0x6E)
INK = RGBColor(0x20, 0x22, 0x2B)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ICE = RGBColor(0xC9, 0xD6, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
CARD_BORDER = RGBColor(0xE1, 0xE5, 0xEC)

GREEN = RGBColor(0x1E, 0x8E, 0x5A)
AMBER = RGBColor(0xC9, 0x82, 0x1B)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY_STATUS = RGBColor(0x8A, 0x90, 0x9C)

HEADING_FONT = "Georgia"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)


class PPTGenerator:

    def __init__(
        self,
        metrics: ProjectMetrics,
        health: HealthResult,
        summary: str,
    ):
        self.metrics = metrics
        self.health = health
        self.summary = summary

    # -----------------------------------------------------------------
    # Generic, defensive accessors
    # -----------------------------------------------------------------

    @staticmethod
    def _get(obj, name, default=None):
        """Read an attribute (or dict key) if present, else fall back.

        The scoring/metrics objects this module receives may carry more
        fields than the ones this presentation layer strictly needs, or
        fewer -- this keeps rendering resilient either way without ever
        touching the upstream models.
        """
        if obj is None:
            return default
        if isinstance(obj, dict):
            return obj.get(name, default)
        return getattr(obj, name, default)

    def _rag_color(self, status=None):
        status = (status or self._get(self.health, "rag_status", "")).strip().lower()
        return {"green": GREEN, "amber": AMBER, "red": RED}.get(status, GREY_STATUS)

    def _score_color(self, value):
        try:
            value = float(value)
        except (TypeError, ValueError):
            return GREY_STATUS
        if value >= 75:
            return GREEN
        if value >= 50:
            return AMBER
        return RED

    def _fmt_score(self, value, suffix="%"):
        try:
            return f"{float(value):.0f}{suffix}"
        except (TypeError, ValueError):
            return "N/A"

    def _normalize_recommendation(self, item):
        """Recommendations may arrive as plain strings or as richer
        objects/dicts (priority, owner, benefit). Handle both without
        assuming a schema the upstream service hasn't guaranteed."""
        if isinstance(item, dict):
            text = item.get("text") or item.get("action") or item.get("recommendation") or str(item)
            priority = item.get("priority", "Medium")
            owner = item.get("owner") or item.get("action_owner") or "Project Team"
            benefit = item.get("benefit") or item.get("expected_benefit") or "Improved execution confidence"
            return text, priority, owner, benefit
        if isinstance(item, str):
            return item, "Medium", "Project Team", "Improved execution confidence"
        # arbitrary object with attributes
        text = self._get(item, "text") or self._get(item, "action") or str(item)
        priority = self._get(item, "priority", "Medium")
        owner = self._get(item, "owner", self._get(item, "action_owner", "Project Team"))
        benefit = self._get(item, "benefit", self._get(item, "expected_benefit", "Improved execution confidence"))
        return text, priority, owner, benefit

    # -----------------------------------------------------------------
    # Low-level drawing helpers
    # -----------------------------------------------------------------

    def _blank_slide(self, prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _fill_background(self, slide, color):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def _textbox(self, slide, x, y, w, h, wrap=True, auto_fit=False):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = wrap
        if auto_fit:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        return box, tf

    def _style_run(self, paragraph, text=None, font=BODY_FONT, size=14, bold=False,
                    italic=False, color=INK, align=None, space_after=None, space_before=None):
        if text is not None:
            paragraph.text = text
        paragraph.font.name = font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        paragraph.font.color.rgb = color
        if align is not None:
            paragraph.alignment = align
        if space_after is not None:
            paragraph.space_after = Pt(space_after)
        if space_before is not None:
            paragraph.space_before = Pt(space_before)
        return paragraph

    def _rounded_card(self, slide, x, y, w, h, fill=LIGHT_GRAY, line_color=CARD_BORDER, line_pt=0.75, radius=None):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        if line_color is None:
            card.line.fill.background()
        else:
            card.line.color.rgb = line_color
            card.line.width = Pt(line_pt)
        card.shadow.inherit = False
        if radius is not None:
            try:
                card.adjustments[0] = radius
            except (IndexError, ValueError):
                pass
        return card

    def _pill(self, slide, x, y, w, h, text, fill, text_color=WHITE, size=16):
        pill = self._rounded_card(slide, x, y, w, h, fill=fill, line_color=None, radius=0.5)
        tf = pill.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        p = tf.paragraphs[0]
        self._style_run(p, text, font=BODY_FONT, size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        tf.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE
        return pill

    def _dot(self, slide, x, y, d, color):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.shadow.inherit = False
        return dot

    def _section_title(self, slide, x, y, w, text, color=NAVY, size=18):
        _, tf = self._textbox(slide, x, y, w, Inches(0.4))
        self._style_run(tf.paragraphs[0], text, font=HEADING_FONT, size=size, bold=True, color=color)
        return tf

    def _bullet_list(self, slide, x, y, w, h, items, color=INK, size=13.5,
                      dot_color=None, space_after=6, auto_fit=True):
        _, tf = self._textbox(slide, x, y, w, h, auto_fit=auto_fit)
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            self._style_run(p, f"•  {item}", font=BODY_FONT, size=size, color=color, space_after=space_after)
        return tf

    def _title_bar(self, slide, kicker, title, page_no, section_color=NAVY):
        """Standard header used on the four white content slides."""
        _, ktf = self._textbox(slide, MARGIN, Inches(0.4), Inches(9), Inches(0.3))
        self._style_run(ktf.paragraphs[0], kicker.upper(), font=BODY_FONT, size=11,
                         bold=True, color=MUTED)

        _, ttf = self._textbox(slide, MARGIN, Inches(0.68), Inches(11), Inches(0.65))
        self._style_run(ttf.paragraphs[0], title, font=HEADING_FONT, size=27, bold=True, color=section_color)

        # Small page marker, top right, consulting-deck style
        _, ptf = self._textbox(slide, Inches(12.2), Inches(0.42), Inches(0.6), Inches(0.3))
        self._style_run(ptf.paragraphs[0], f"{page_no:02d}", font=BODY_FONT, size=11,
                         color=MUTED, align=PP_ALIGN.RIGHT)

    def _footer(self, slide, page_no, dark=False):
        color = ICE if dark else MUTED
        _, tf = self._textbox(slide, MARGIN, Inches(7.12), Inches(9), Inches(0.3))
        self._style_run(
            tf.paragraphs[0],
            "CONFIDENTIAL  |  Project Health Reporting Agent — Executive Report",
            font=BODY_FONT, size=9, color=color,
        )
        _, ptf = self._textbox(slide, Inches(12.2), Inches(7.12), Inches(0.6), Inches(0.3))
        self._style_run(ptf.paragraphs[0], f"{page_no:02d}", font=BODY_FONT, size=9,
                         color=color, align=PP_ALIGN.RIGHT)

    def _kpi_card(self, slide, x, y, w, h, label, value, score_for_color=None):
        card = self._rounded_card(slide, x, y, w, h, fill=WHITE, line_color=CARD_BORDER, radius=0.10)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.14)

        dot_color = self._score_color(score_for_color) if score_for_color is not None else NAVY
        self._dot(slide, x + Inches(0.15), y + Inches(0.16), Inches(0.14), dot_color)

        label_p = tf.paragraphs[0]
        self._style_run(label_p, "   " + label, font=BODY_FONT, size=12.5, bold=True, color=MUTED)

        value_p = tf.add_paragraph()
        self._style_run(value_p, value, font=HEADING_FONT, size=30, bold=True, color=NAVY, space_before=6)

    def _bar_chart(self, slide, x, y, w, h, categories, values, title="Health Signal Breakdown"):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("Score", values)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = False
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        title_p = chart.chart_title.text_frame.paragraphs[0]
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.name = BODY_FONT
        title_p.font.color.rgb = NAVY

        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = "0"
        plot.data_labels.number_format_is_linked = False
        plot.data_labels.font.size = Pt(11)
        plot.data_labels.font.bold = True
        plot.data_labels.font.name = BODY_FONT
        plot.data_labels.font.color.rgb = INK
        plot.gap_width = 60

        series = plot.series[0]
        series_values = list(series.values)
        for idx, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self._score_color(series_values[idx])
            point.format.line.fill.background()

        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(11)
        cat_axis.tick_labels.font.name = BODY_FONT
        cat_axis.tick_labels.font.color.rgb = INK
        cat_axis.format.line.color.rgb = CARD_BORDER

        val_axis = chart.value_axis
        val_axis.minimum_scale = 0
        val_axis.maximum_scale = 100
        val_axis.tick_labels.font.size = Pt(10)
        val_axis.tick_labels.font.name = BODY_FONT
        val_axis.tick_labels.font.color.rgb = MUTED
        val_axis.major_gridlines.format.line.color.rgb = CARD_BORDER
        val_axis.major_gridlines.format.line.width = Pt(0.5)
        val_axis.format.line.fill.background()

        return graphic_frame

    # -----------------------------------------------------------------
    # Slides
    # -----------------------------------------------------------------

    def _slide_cover(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, NAVY)

        # Soft decorative depth without a literal accent stripe/bar.
        glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.6), Inches(-2.0), Inches(6.0), Inches(6.0))
        glow.fill.solid()
        glow.fill.fore_color.rgb = NAVY_SOFT
        glow.line.fill.background()
        glow.shadow.inherit = False

        _, ktf = self._textbox(slide, MARGIN, Inches(0.7), Inches(9), Inches(0.35))
        self._style_run(ktf.paragraphs[0], "MONTHLY PROJECT HEALTH REVIEW", font=BODY_FONT,
                         size=13, bold=True, color=ICE)

        _, ttf = self._textbox(slide, MARGIN, Inches(1.05), Inches(9.4), Inches(1.3))
        self._style_run(ttf.paragraphs[0], "Project Health Executive Report",
                         font=HEADING_FONT, size=34, bold=True, color=WHITE)

        project_name = self._get(self.metrics, "project_name", "Unnamed Project")
        project_manager = self._get(self.metrics, "project_manager", "Unassigned")
        generated_date = datetime.now().strftime("%d %B %Y")

        _, stf = self._textbox(slide, MARGIN, Inches(2.15), Inches(8.5), Inches(0.5))
        self._style_run(stf.paragraphs[0], project_name, font=BODY_FONT, size=20, bold=True, color=ICE)

        _, mtf = self._textbox(slide, MARGIN, Inches(2.65), Inches(9), Inches(0.4))
        self._style_run(
            mtf.paragraphs[0],
            f"Project Manager: {project_manager}      |      Generated: {generated_date}",
            font=BODY_FONT, size=13, color=RGBColor(0x9A, 0xA6, 0xD6),
        )

        # Large RAG badge, top right
        rag_status = self._get(self.health, "rag_status", "N/A")
        self._pill(slide, Inches(10.3), Inches(0.7), Inches(2.4), Inches(1.05),
                   str(rag_status).upper(), self._rag_color(), size=26)

        # Stat callouts: Overall Score / Confidence
        overall_score = self._get(self.health, "overall_score", None)
        confidence = self._get(self.health, "confidence_score", None)

        for i, (label, val, suffix) in enumerate([
            ("OVERALL HEALTH SCORE", overall_score, "/100"),
            ("CONFIDENCE", confidence, "%"),
        ]):
            cx = MARGIN + Emu(int(i * Inches(4.6)))
            card = self._rounded_card(slide, cx, Inches(3.5), Inches(4.2), Inches(2.15),
                                       fill=NAVY_CARD, line_color=None, radius=0.08)
            tf = card.text_frame
            tf.margin_left = Inches(0.25)
            tf.margin_top = Inches(0.2)
            p = tf.paragraphs[0]
            self._style_run(p, label, font=BODY_FONT, size=12, bold=True, color=ICE)
            vp = tf.add_paragraph()
            value_text = f"{val:.1f}{suffix}" if isinstance(val, (int, float)) else "N/A"
            self._style_run(vp, value_text, font=HEADING_FONT, size=44, bold=True, color=WHITE, space_before=6)

        self._footer(slide, page_no=1, dark=True)
        return slide

    def _slide_dashboard(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Executive Dashboard", "Project Performance at a Glance", page_no=2)

        completion = self._get(self.health, "completion_score", 0)
        schedule = self._get(self.health, "schedule_score", 0)
        milestone = self._get(self.health, "milestone_score", 0)
        blockers = self._get(self.health, "blocker_score", 0)

        kpis = [
            ("Completion", completion, self._fmt_score(completion)),
            ("Schedule", schedule, self._fmt_score(schedule)),
            ("Milestones", milestone, self._fmt_score(milestone)),
            ("Blockers", blockers, self._fmt_score(blockers, suffix="")),
        ]
        card_w = Inches(2.85)
        card_h = Inches(1.35)
        gap = Inches(0.25)
        start_x = MARGIN
        for i, (label, raw, display) in enumerate(kpis):
            x = start_x + Emu(int(i * (card_w + gap)))
            self._kpi_card(slide, x, Inches(1.55), card_w, card_h, label, display, score_for_color=raw)

        # Chart, left column
        self._bar_chart(
            slide, MARGIN, Inches(3.35), Inches(7.6), Inches(3.55),
            categories=["Completion", "Schedule", "Milestones", "Blockers"],
            values=[completion, schedule, milestone, blockers],
        )

        # Executive Insights, right column
        insights_x = Inches(8.55)
        self._rounded_card(slide, insights_x, Inches(3.35), Inches(4.25), Inches(3.55),
                            fill=LIGHT_GRAY, line_color=CARD_BORDER, radius=0.06)
        self._section_title(slide, insights_x + Inches(0.3), Inches(3.55), Inches(3.7),
                             "Executive Insights", size=16)
        self._bullet_list(
            slide, insights_x + Inches(0.3), Inches(4.05), Inches(3.7), Inches(2.7),
            items=[
                "Overall execution remains healthy with strong schedule adherence.",
                "Milestone completion requires closer monitoring.",
                "Critical tasks should be prioritized to avoid future delays.",
            ],
            size=13, space_after=10,
        )

        self._footer(slide, page_no=2)
        return slide

    def _slide_risks(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Risk Assessment", "Key Risks & Business Impact", page_no=3)

        rag_status = str(self._get(self.health, "rag_status", "N/A")).upper()
        self._pill(slide, Inches(10.3), Inches(0.55), Inches(2.4), Inches(0.55),
                   f"RISK LEVEL: {rag_status}", self._rag_color(), size=13)

        reasons = self._get(self.health, "reasons", []) or ["No specific risks captured for this cycle."]
        impacts = [
            "Potential milestone delay",
            "Higher delivery risk",
            "Need for closer executive monitoring",
        ]

        col_w = Inches(5.9)
        left_x = MARGIN
        right_x = Inches(6.85)
        top_y = Inches(1.55)
        col_h = Inches(4.9)

        left_card = self._rounded_card(slide, left_x, top_y, col_w, col_h, fill=LIGHT_GRAY,
                                        line_color=CARD_BORDER, radius=0.05)
        self._dot(slide, left_x + Inches(0.3), top_y + Inches(0.32), Inches(0.16), self._rag_color())
        self._section_title(slide, left_x + Inches(0.55), top_y + Inches(0.22), col_w - Inches(0.8),
                             "Critical Risks", size=17)
        self._bullet_list(slide, left_x + Inches(0.3), top_y + Inches(0.85), col_w - Inches(0.6),
                           col_h - Inches(1.1), items=reasons, size=13.5,
                           color=INK, space_after=10)

        right_card = self._rounded_card(slide, right_x, top_y, col_w, col_h, fill=LIGHT_GRAY,
                                         line_color=CARD_BORDER, radius=0.05)
        self._dot(slide, right_x + Inches(0.3), top_y + Inches(0.32), Inches(0.16), NAVY)
        self._section_title(slide, right_x + Inches(0.55), top_y + Inches(0.22), col_w - Inches(0.8),
                             "Business Impact", size=17)
        self._bullet_list(slide, right_x + Inches(0.3), top_y + Inches(0.85), col_w - Inches(0.6),
                           col_h - Inches(1.1), items=impacts, size=13.5,
                           color=INK, space_after=10)

        self._footer(slide, page_no=3)
        return slide

    def _slide_recommendations(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Recommendations", "Priority Actions & Expected Impact", page_no=4)

        raw_recs = self._get(self.health, "recommendations", []) or ["No specific recommendations for this cycle."]
        priority_color = {"high": RED, "medium": AMBER, "low": GREEN}

        top_y = Inches(1.55)
        row_h = Inches(1.05)
        gap = Inches(0.18)
        max_rows = 4
        for i, item in enumerate(raw_recs[:max_rows]):
            text, priority, owner, benefit = self._normalize_recommendation(item)
            y = top_y + Emu(int(i * (row_h + gap)))

            card = self._rounded_card(slide, MARGIN, y, Inches(12.23), row_h, fill=LIGHT_GRAY,
                                       line_color=CARD_BORDER, radius=0.12)

            # Index badge
            idx_d = Inches(0.55)
            idx_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, MARGIN + Inches(0.25),
                                                y + Inches(0.25), idx_d, idx_d)
            idx_shape.fill.solid()
            idx_shape.fill.fore_color.rgb = NAVY
            idx_shape.line.fill.background()
            idx_shape.shadow.inherit = False
            itf = idx_shape.text_frame
            itf.word_wrap = True
            self._style_run(itf.paragraphs[0], str(i + 1), font=HEADING_FONT, size=18,
                             bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            _, ttf = self._textbox(slide, MARGIN + Inches(1.05), y + Inches(0.10), Inches(7.7), Inches(0.42))
            self._style_run(ttf.paragraphs[0], text, font=BODY_FONT, size=13, bold=True, color=INK)

            _, benefit_tf = self._textbox(slide, MARGIN + Inches(1.05), y + Inches(0.60), Inches(7.7), Inches(0.4))
            self._style_run(benefit_tf.paragraphs[0], f"Expected benefit: {benefit}", font=BODY_FONT,
                             size=11, italic=True, color=MUTED)

            self._pill(slide, Inches(9.0), y + Inches(0.16), Inches(1.55), Inches(0.35),
                       str(priority).upper(), priority_color.get(str(priority).lower(), AMBER), size=11)

            _, owner_tf = self._textbox(slide, Inches(10.65), y + Inches(0.10), Inches(1.75), Inches(0.85))
            self._style_run(owner_tf.paragraphs[0], "OWNER", font=BODY_FONT, size=9, bold=True, color=MUTED)
            op = owner_tf.add_paragraph()
            self._style_run(op, owner, font=BODY_FONT, size=12, bold=True, color=NAVY, space_before=2)

        self._footer(slide, page_no=4)
        return slide

    def _slide_executive_summary(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, WHITE)
        self._title_bar(slide, "Executive Summary", "From Data to Decision", page_no=5)

        rag_status = str(self._get(self.health, "rag_status", "N/A")).upper()
        overall_score = self._get(self.health, "overall_score", "N/A")
        confidence = self._get(self.health, "confidence_score", "N/A")
        reasons = self._get(self.health, "reasons", []) or ["No specific risks captured."]
        recommendations = self._get(self.health, "recommendations", []) or ["No specific recommendations."]
        impacts = ["Potential milestone delay", "Higher delivery risk", "Need for closer executive monitoring"]

        # Situation — full-width narrative card (the AI-generated summary)
        sit_h = Inches(1.55)
        self._rounded_card(slide, MARGIN, Inches(1.5), Inches(12.23), sit_h, fill=LIGHT_GRAY,
                            line_color=CARD_BORDER, radius=0.06)
        self._section_title(slide, MARGIN + Inches(0.3), Inches(1.62), Inches(11.6), "Situation", size=15)
        _, sit_tf = self._textbox(slide, MARGIN + Inches(0.3), Inches(2.0), Inches(11.6), Inches(1.0), auto_fit=True)
        self._style_run(sit_tf.paragraphs[0], self.summary, font=BODY_FONT, size=12.5, color=INK)

        # Four supporting sections in a 2x2 grid, each its own labeled card
        sections = [
            ("Current Status", [f"{rag_status} — Overall Score {overall_score}/100, Confidence {confidence}%"], NAVY),
            ("Key Risks", reasons[:3], self._rag_color()),
            ("Business Impact", impacts, NAVY),
            ("Recommendations", [self._normalize_recommendation(r)[0] for r in recommendations[:3]], GREEN),
        ]
        grid_top = Inches(3.3)
        cell_w = Inches(5.95)
        cell_h = Inches(1.75)
        gap_x = Inches(0.3)
        gap_y = Inches(0.2)
        for i, (label, items, accent) in enumerate(sections):
            col = i % 2
            row = i // 2
            x = MARGIN + Emu(int(col * (cell_w + gap_x)))
            y = grid_top + Emu(int(row * (cell_h + gap_y)))
            self._rounded_card(slide, x, y, cell_w, cell_h, fill=WHITE, line_color=CARD_BORDER, radius=0.06)
            self._dot(slide, x + Inches(0.25), y + Inches(0.22), Inches(0.14), accent)
            self._section_title(slide, x + Inches(0.5), y + Inches(0.12), cell_w - Inches(0.7), label, size=14)
            self._bullet_list(slide, x + Inches(0.25), y + Inches(0.58), cell_w - Inches(0.5),
                               cell_h - Inches(0.7), items=items, size=11.5, space_after=4)

        self._footer(slide, page_no=5)
        return slide

    def _slide_snapshot(self, prs):
        slide = self._blank_slide(prs)
        self._fill_background(slide, NAVY)

        _, ktf = self._textbox(slide, MARGIN, Inches(0.5), Inches(9), Inches(0.3))
        self._style_run(ktf.paragraphs[0], "PROJECT SNAPSHOT", font=BODY_FONT, size=12, bold=True, color=ICE)

        _, ttf = self._textbox(slide, MARGIN, Inches(0.8), Inches(10), Inches(0.65))
        self._style_run(ttf.paragraphs[0], "Snapshot & Conclusion", font=HEADING_FONT, size=30, bold=True, color=WHITE)

        project_name = self._get(self.metrics, "project_name", "Unnamed Project")
        project_manager = self._get(self.metrics, "project_manager", "Unassigned")

        _, idf = self._textbox(slide, MARGIN, Inches(1.55), Inches(6), Inches(0.6))
        self._style_run(idf.paragraphs[0], project_name, font=BODY_FONT, size=16, bold=True, color=ICE)
        p2 = idf.add_paragraph()
        self._style_run(p2, f"Managed by {project_manager}", font=BODY_FONT, size=12,
                         color=RGBColor(0x9A, 0xA6, 0xD6), space_before=2)

        stats = [
            ("Total Tasks", self._get(self.metrics, "total_tasks", "N/A")),
            ("Completed", self._get(self.metrics, "completed_tasks", "N/A")),
            ("Milestones", self._get(self.metrics, "milestone_count", "N/A")),
            ("Critical Tasks", self._get(self.metrics, "critical_task_count", "N/A")),
            ("Overall Score", self._fmt_score(self._get(self.health, "overall_score", None), suffix="/100")),
            ("Confidence", self._fmt_score(self._get(self.health, "confidence_score", None))),
        ]
        chip_w = Inches(1.92)
        chip_h = Inches(1.05)
        gap = Inches(0.15)
        start_x = MARGIN
        start_y = Inches(2.35)
        for i, (label, value) in enumerate(stats):
            col = i % 3
            row = i // 3
            x = start_x + Emu(int(col * (chip_w + gap)))
            y = start_y + Emu(int(row * (chip_h + gap)))
            chip = self._rounded_card(slide, x, y, chip_w, chip_h, fill=NAVY_CARD, line_color=None, radius=0.1)
            tf = chip.text_frame
            tf.margin_left = Inches(0.12)
            tf.margin_top = Inches(0.1)
            self._style_run(tf.paragraphs[0], label, font=BODY_FONT, size=10.5, bold=True, color=ICE)
            vp = tf.add_paragraph()
            self._style_run(vp, str(value), font=HEADING_FONT, size=20, bold=True, color=WHITE, space_before=4)

        # Final conclusion — light card floating on the dark canvas for contrast
        concl_x = Inches(7.6)
        concl_y = Inches(1.55)
        concl_w = Inches(5.15)
        concl_h = Inches(4.35)
        self._rounded_card(slide, concl_x, concl_y, concl_w, concl_h, fill=WHITE, line_color=None, radius=0.05)
        self._section_title(slide, concl_x + Inches(0.3), concl_y + Inches(0.22), concl_w - Inches(0.6),"Final Conclusion", color=NAVY, size=17)
        self._bullet_list(
            slide, concl_x + Inches(0.3), concl_y + Inches(0.75), concl_w - Inches(0.6), Inches(2.0),
            items=[
                "Automated Excel Parsing",
                "Intelligent KPI Extraction",
                "Automated Health Scoring Engine",
                "AI Executive Summary using Groq Llama 3.3",
                "Automatic PowerPoint Generation",
            ],
            color=INK, size=12.5, space_after=6,
        )
        _, note_tf = self._textbox(slide, concl_x + Inches(0.3), concl_y + Inches(3.15),
                                    concl_w - Inches(0.6), Inches(1.0), auto_fit=True)
        self._style_run(
            note_tf.paragraphs[0],
            "This solution demonstrates an end-to-end AI-powered Project Health "
            "Reporting workflow suitable for executive decision making.",
            font=BODY_FONT, size=11.5, italic=True, color=MUTED,
        )

        _, thank_tf = self._textbox(slide, MARGIN, Inches(6.15), Inches(7), Inches(0.6))
        self._style_run(thank_tf.paragraphs[0], "Thank You", font=HEADING_FONT, size=26, bold=True, color=WHITE)

        self._footer(slide, page_no=6, dark=True)
        return slide

    # -----------------------------------------------------------------
    # Entry point
    # -----------------------------------------------------------------

    def generate(self, output_path):
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        self._slide_cover(prs)
        self._slide_dashboard(prs)
        self._slide_risks(prs)
        self._slide_recommendations(prs)
        self._slide_executive_summary(prs)
        self._slide_snapshot(prs)

        prs.save(output_path)